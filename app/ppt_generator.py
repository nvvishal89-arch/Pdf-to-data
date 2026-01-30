"""
Phase 3: Auto PPT Generator (spec §9).
Slides: Project Summary → Product Overview → Product Renders → Technical Drawings → Lifecycle → Delivery.
"""
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt

from app.schema import SQStructuredData


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def _add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.add_paragraph() if i else body.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(14)


def _add_product_slide(prs: Presentation, product, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left = Inches(0.5)
    top = Inches(0.5)
    w, h = Inches(9), Inches(1)
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.text = f"Product {index + 1}/{total}: {product.name}"
    for p in tf.paragraphs:
        p.font.size = Pt(24)
        p.font.bold = True
    # Description
    tb2 = slide.shapes.add_textbox(left, Inches(1.6), w, Inches(2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p0 = tf2.paragraphs[0]
    p0.text = (product.description or "")[:300] + ("…" if len(product.description or "") > 300 else "")
    p0.font.size = Pt(12)
    if product.dimensions:
        p2 = tf2.add_paragraph()
        p2.text = f"Dimensions: {product.dimensions}"
        p2.font.size = Pt(11)
    p3 = tf2.add_paragraph()
    p3.text = f"Qty: {product.qty}  |  Rate: ₹{product.unit_price:,.0f}  |  Amount: ₹{product.amount:,.0f}"
    p3.font.size = Pt(11)
    # Optional: add first product image if base64 available
    if product.images:
        try:
            import base64
            from pptx.util import Emu
            img_data = base64.b64decode(product.images[0])
            pic = slide.shapes.add_picture(
                BytesIO(img_data), Inches(0.5), Inches(2.8), width=Inches(2.5)
            )
        except Exception:
            pass


def generate_ppt(data: SQStructuredData) -> bytes:
    """Generate PowerPoint from SQ structured data. Returns .pptx file bytes."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    proj = data.project
    # 1. Project Summary
    _add_title_slide(
        prs,
        "Project Summary",
        f"{proj.project_name or 'SQ'}  |  Client: {proj.client_name or '—'}  |  Date: {proj.date or '—'}"
    )
    # 2. Product Overview
    bullets = [f"{p.sr_no}. {p.name} – Qty {p.qty}, ₹{p.amount:,.0f}" for p in data.products[:15]]
    if len(data.products) > 15:
        bullets.append(f"... and {len(data.products) - 15} more")
    _add_content_slide(prs, "Product Overview", bullets or ["No products"])
    # 3. Product Render Slides
    for i, p in enumerate(data.products):
        _add_product_slide(prs, p, i, len(data.products))
    # 4. Technical Drawings (placeholder)
    _add_content_slide(prs, "Technical Drawings", ["Per product drawings (Phase 2 output)."])
    # 5. Manufacturing Lifecycle (placeholder)
    _add_content_slide(prs, "Manufacturing Lifecycle", [
        "Machining → Carpentry → Metal → Assembly → Upholstery → Paint → Final Assembly → Packaging → Dispatch",
        "(See SOW for per-product stages.)",
    ])
    # 6. Delivery Timeline (placeholder)
    _add_content_slide(prs, "Delivery Timeline", ["TBD – link to Gantt (Phase 4)."])

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
